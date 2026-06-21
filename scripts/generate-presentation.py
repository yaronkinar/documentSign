# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

NAVY = RGBColor(0x0F, 0x2A, 0x4A)
BLUE = RGBColor(0x2E, 0x6F, 0xE0)
LIGHT_BLUE = RGBColor(0xEA, 0xF1, 0xFD)
GRAY = RGBColor(0x4A, 0x4A, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x1F, 0x9D, 0x55)
ORANGE = RGBColor(0xE0, 0x7B, 0x2E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]


def rtl_run(run):
    rPr = run._r.get_or_add_rPr()
    rPr.set('lang', 'he-IL')


def set_rtl_para(p):
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    pPr.set('rtl', '1')
    p.alignment = PP_ALIGN.RIGHT


def add_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_text(slide, left, top, width, height, text, size=20, bold=False,
             color=NAVY, align=PP_ALIGN.RIGHT, font="Arial", anchor=None,
             line_spacing=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        set_rtl_para(p)
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
        rtl_run(r)
    return tb


def add_bullets(slide, left, top, width, height, items, size=15, color=GRAY,
                 bold_first=False, space_after=8, marker="•  "):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.RIGHT
        set_rtl_para(p)
        p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = marker + item
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Arial"
        rtl_run(r)
    return tb


def add_footer(slide, page_num, total):
    add_text(slide, Inches(0.4), Inches(7.1), Inches(3), Inches(0.3),
              "DocFlows", size=11, color=RGBColor(0xAA,0xAA,0xAA), align=PP_ALIGN.LEFT)
    add_text(slide, SW - Inches(3.4), Inches(7.1), Inches(3), Inches(0.3),
              f"{page_num} / {total}", size=11, color=RGBColor(0xAA,0xAA,0xAA), align=PP_ALIGN.RIGHT)


def section_header(slide, kicker, title, sub=None):
    add_text(slide, Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.5),
              kicker, size=16, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.7), Inches(0.85), Inches(11.9), Inches(0.9),
              title, size=30, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.55), Inches(1.4), Pt(4))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background(); bar.shadow.inherit = False
    if sub:
        add_text(slide, Inches(0.7), Inches(1.7), Inches(11.9), Inches(0.5),
                  sub, size=15, color=GRAY, align=PP_ALIGN.RIGHT)


def feature_card(slide, left, top, width, height, title, desc, accent=BLUE):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.06
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xE2, 0xE6, 0xEC); card.line.width = Pt(1)
    card.shadow.inherit = False
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + width - Inches(0.08), top, Inches(0.08), height)
    accent_bar.fill.solid(); accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background(); accent_bar.shadow.inherit = False
    add_text(slide, left + Inches(0.25), top + Inches(0.15), width - Inches(0.55), Inches(0.5),
              title, size=15, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
    add_text(slide, left + Inches(0.25), top + Inches(0.62), width - Inches(0.55), height - Inches(0.75),
              desc, size=11.5, color=GRAY, align=PP_ALIGN.RIGHT)


TOTAL = 14
n = 0

# ---------- Slide 1: Title ----------
s = prs.slides.add_slide(blank)
add_bg(s, NAVY)
add_text(s, Inches(1), Inches(2.5), Inches(11.3), Inches(1.2),
          "DocFlows", size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.7), Inches(11.3), Inches(0.8),
          "פלטפורמת חתימה דיגיטלית וניהול חוזים חכמה", size=24, color=RGBColor(0xCF,0xDF,0xFF), align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.4), Inches(11.3), Inches(0.6),
          "מבוססת AI • תומכת בעברית ובטפסים ממשלתיים ועירוניים", size=16, color=RGBColor(0x9C,0xB6,0xE8), align=PP_ALIGN.CENTER)
n += 1; add_footer(s, n, TOTAL)

# ---------- Slide 2: Agenda ----------
s = prs.slides.add_slide(blank)
add_bg(s)
section_header(s, "סדר היום", "מה נסקור במצגת")
agenda = [
    "האתגר העסקי – מה היה לפני DocFlows",
    "סקירת המוצר – זרימת העבודה המרכזית",
    "חתימה דיגיטלית וניהול תהליכי אישור",
    "אוטומציה מבוססת AI – הזמן שהיא חוסכת",
    "תמיכה בטפסי הקנסות / טפסים עירוניים בעברית",
    "תבניות, שיתוף פעולה ותקשורת",
    "אבטחה, מעקב ותאימות (Audit Trail)",
    "סיכום היתרונות והערך העסקי",
]
add_bullets(s, Inches(1.2), Inches(2.1), Inches(10.5), Inches(4.5), agenda, size=18, space_after=14)
n += 1; add_footer(s, n, TOTAL)

# ---------- Slide 3: The Problem ----------
s = prs.slides.add_slide(blank)
add_bg(s)
section_header(s, "האתגר", "ניהול חוזים וחתימות - הכאב האמיתי")
problems = [
    ("הדפסה, חתימה וסריקה", "תהליכים ידניים שגוזלים זמן וכרוכים בנסיעות ופגישות פיזיות"),
    ("איתור חוזרני של פרטי חוזה", "הקלדה חזרתית של מספרי חוזה, תאריכים, סכומים וספקים מתוך מסמכים"),
    ("היעדר מעקב ותיעוד", "קושי לדעת מי חתם, מתי, ובאיזה שלב נמצא האישור"),
    ("טפסים עירוניים מורכבים", "טופס הקנסות ודומיו דורשים דיוק רב והתאמה לדרישות רגולטוריות"),
]
top = Inches(2.0)
for i, (t, d) in enumerate(problems):
    feature_card(s, Inches(0.7), top + i*Inches(1.15), Inches(11.9), Inches(1.0), t, d, accent=ORANGE)
n += 1; add_footer(s, n, TOTAL)

# ---------- Slide 4: Solution overview ----------
s = prs.slides.add_slide(blank)
add_bg(s)
section_header(s, "הפתרון", "DocFlows - כל תהליך החוזה במקום אחד")
add_text(s, Inches(0.7), Inches(2.0), Inches(11.9), Inches(0.6),
          "העלאה → חילוץ נתונים אוטומטי ב-AI → חתימה ואישור → מעקב מלא → ארכיון", size=17, color=BLUE, bold=True)
steps = ["העלאת מסמך\n(PDF / Word)", "AI מחלץ שדות\nונתוני חוזה", "חתימה ואישור\nרב-שלבי", "מעקב, תגובות\nותיעוד מלא"]
colors = [BLUE, GREEN, ORANGE, NAVY]
w = Inches(2.7); gap = Inches(0.25)
total_w = w*4 + gap*3
start_x = (SW - total_w)//2
for i, (t, c) in enumerate(zip(steps, colors)):
    x = start_x + i*(w+gap)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.0), w, Inches(1.6))
    box.fill.solid(); box.fill.fore_color.rgb = c; box.line.fill.background(); box.shadow.inherit = False
    box.adjustments[0] = 0.1
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    for j, line in enumerate(t.split("\n")):
        pp = p if j == 0 else tf.add_paragraph()
        pp.alignment = PP_ALIGN.CENTER
        r = pp.add_run(); r.text = line; r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name="Arial"
        rtl_run(r)
    if i < 3:
        arrow = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, x - gap - Inches(0.05), Inches(3.6), gap, Inches(0.4))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = RGBColor(0xC8,0xC8,0xC8); arrow.line.fill.background(); arrow.shadow.inherit=False
add_text(s, Inches(0.7), Inches(5.1), Inches(11.9), Inches(1.5),
          "התוצאה: זמן הקמת חוזה מתקצר מימים לדקות, שגיאות הקלדה יורדות, ולכל פעולה יש תיעוד מלא לצורכי ביקורת ותאימות משפטית.",
          size=16, color=GRAY)
n += 1; add_footer(s, n, TOTAL)

# ---------- Slide 5: Core signing features ----------
s = prs.slides.add_slide(blank)
add_bg(s)
section_header(s, "תכונת ליבה", "העלאה והמרת מסמכים")
cards = [
    ("תמיכה רב-פורמטית", "העלאת PDF, DOC ו-DOCX; המרה אוטומטית של קבצי Word ל-PDF ללא צורך בכלים חיצוניים"),
    ("שמירה על נאמנות המקור", "המסמך המומר שומר על העימוד, הפונטים והעיצוב המקוריים"),
    ("API להמרה עצמאית", "נקודת קצה ייעודית להמרת Word ל-PDF בדרישה, לשילוב בתהליכים נוספים"),
]
top = Inches(2.1)
for i, (t, d) in enumerate(cards):
    feature_card(s, Inches(0.7), top + i*Inches(1.5), Inches(11.9), Inches(1.3), t, d)
n += 1; add_footer(s, n, TOTAL)

# ---------- Slide 6: Signature & workflow ----------
s = prs.slides.add_slide(blank)
add_bg(s)
section_header(s, "תכונת ליבה", "חתימה דיגיטלית ותהליכי אישור")
left_items = [
    "חתימה בשלוש דרכים: ציור, הקלדה והעלאת תמונה",
    "מיקום מדויק (פיקסל) של שדות חתימה על העמוד",
    "חותמים אורחים ללא צורך בהרשמה, באמצעות קישור מאובטח שתקף ל-72 שעות",
    "גלריית חתימות לשימוש חוזר",
]
right_items = [
    "תהליכי אישור רב-שלביים - רציפים או מקבילים",
    "ארבעה סוגי שלבים: חתימה, סקירה, אישור, התראה",
    "מעקב סטטוס מסמך: טיוטה → בהמתנה לחתימה/אישור → אושר",
    "אפשרות \"דילוג\" לחותמים לא-חובה",
]
add_text(s, Inches(7.1), Inches(2.05), Inches(5.5), Inches(0.5), "חתימה וחותמים", size=18, bold=True, color=BLUE)
add_bullets(s, Inches(6.9), Inches(2.6), Inches(5.7), Inches(3.5), left_items, size=14.5)
add_text(s, Inches(0.7), Inches(2.05), Inches(5.5), Inches(0.5), "ניהול תהליכי אישור", size=18, bold=True, color=BLUE)
add_bullets(s, Inches(0.5), Inches(2.6), Inches(5.7), Inches(3.5), right_items, size=14.5)
n += 1; add_footer(s, n, TOTAL)

# ---------- Slide 7: AI section intro ----------
s = prs.slides.add_slide(blank)
add_bg(s, NAVY)
add_text(s, Inches(1), Inches(2.6), Inches(11.3), Inches(1.0),
          "אוטומציה מבוססת AI", size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(3.7), Inches(10.3), Inches(1.2),
          "הליבה הטכנולוגית של DocFlows: מנוע AI שקורא את החוזה בעצמו -\nמסכם, מאתר חותמים, מאתר שדות חתימה וממלא נתונים אוטומטית",
          size=18, color=RGBColor(0xCF,0xDF,0xFF), align=PP_ALIGN.CENTER)
n += 1; add_footer(s, n, TOTAL)

# ---------- Slide 8: AI features grid ----------
s = prs.slides.add_slide(blank)
add_bg(s)
section_header(s, "AI", "ארבע יכולות AI שחוסכות שעות עבודה")
ai_cards = [
    ("סיכום מסמך אוטומטי", "AI מסכם את החוזה ב-5 משפטים בהתבסס על הכותרת, השדות והחותמים - בעברית ובאנגלית. חוסך זמן קריאה לחותמים."),
    ("איתור חותמים אוטומטי", "המערכת מזהה מתוך טקסט המסמך את תפקידי החותמים והמאשרים, ומבטלת איתור ידני."),
    ("חילוץ שדות חתימה מ-PDF", "AI חזותי (Vision) מאתר שורות חתימה, שדות אישור ומקומות מילוי בקובץ סרוק או ישן, כולל גיבוי טקסטואלי לסקירה ירודה."),
    ("מילוי אוטומטי של נתוני חוזה", "התכונה החדשה ביותר: AI מחלץ מספר חוזה, תאריכים, פרטי ספק וסכומים ומציע מילוי מוקדם של השדות - ומבטל הקלדה חזרתית ושגיאות אדם."),
]
positions = [(Inches(0.7), Inches(2.0)), (Inches(6.85), Inches(2.0)), (Inches(0.7), Inches(4.0)), (Inches(6.85), Inches(4.0))]
for (t, d), (x, y) in zip(ai_cards, positions):
    feature_card(s, x, y, Inches(5.75), Inches(1.8), t, d, accent=GREEN)
n += 1; add_footer(s, n, TOTAL)

# ---------- Slide 9: Haknasot ----------
s = prs.slides.add_slide(blank)
add_bg(s)
section_header(s, "ייעודי לישראל", "תמיכה מובנית בטופס הקנסות העירוני")
items = [
    "תבנית מובנית לטופס אישור עירוני בן 4 עמודים, התואמת לדרישות רכש עירוניות בישראל",
    "תמיכה ב-30+ שדות חובה, תיבות סיווג (סוג חוזה, סוג קרבה) ושורות אישור דינמיות לפי תפקיד",
    "ממשק להתאמת מיקום שדות, כדי להתאים בדיוק לחלל הריק בטופס המקורי",
    "הפקת PDF סופי וחתום עם כל הנתונים והחתימות משולבים בו",
]
add_bullets(s, Inches(0.9), Inches(2.1), Inches(11.3), Inches(3.0), items, size=16, space_after=14)
val = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(5.3), Inches(11.3), Inches(1.1))
val.fill.solid(); val.fill.fore_color.rgb = LIGHT_BLUE; val.line.fill.background(); val.shadow.inherit=False
add_text(s, Inches(1.2), Inches(5.45), Inches(10.7), Inches(0.8),
          "הערך העסקי: עמידה רגולטורית ללא צורך במומחה טפסים, ביטול שגיאות מילוי, וחיסכון משמעותי בזמן הכנת אישורי רכש עירוניים.",
          size=14.5, color=NAVY, bold=True)
n += 1; add_footer(s, n, TOTAL)

# ---------- Slide 10: Templates & reuse ----------
s = prs.slides.add_slide(blank)
add_bg(s)
section_header(s, "תכונה", "תבניות ושימוש חוזר")
cards = [
    ("תבניות PDF", "שמירת מסמך שהושלם כתבנית, כולל מיקומי שדות החתימה - ויצירת מסמכים חדשים בלחיצה אחת מתוך התבנית"),
    ("תבניות טפסים", "צירוף טפסים (כגון הקנסות) למסמך אחרי ההעלאה, עם קטלוג שדות מוגדר מראש וניהול גמיש של שדות"),
]
top = Inches(2.1)
for i, (t, d) in enumerate(cards):
    feature_card(s, Inches(0.7), top + i*Inches(1.6), Inches(11.9), Inches(1.4), t, d, accent=BLUE)
add_text(s, Inches(0.7), Inches(5.6), Inches(11.9), Inches(0.8),
          "הערך העסקי: עבור חוזים חזרתיים (כגון הסכמי שכירות או רכש סטנדרטיים) - זמן ההקמה יורד מדקות ארוכות לשניות.",
          size=14.5, color=GRAY)
n += 1; add_footer(s, n, TOTAL)

# ---------- Slide 11: Collaboration ----------
s = prs.slides.add_slide(blank)
add_bg(s)
section_header(s, "תכונה", "שיתוף פעולה, תגובות והתראות")
left = [
    "תגובות מעוגנות למקום מדויק בעמוד ה-PDF",
    "שרשור תגובות ותשובות (Threading)",
    "ארבעה סוגי תגובה: כלליות, הערת עיצוב, סיבת דחייה, הערת אישור",
    "אזכור (@) חותמים שמפעיל התראת מייל",
    "תהליך \"סגירת\" תגובה לפני אישור סופי של המסמך",
    "עדכונים בזמן אמת (WebSocket) - כל המשתמשים רואים שינויים מיידית",
]
right = [
    "מיילי הזמנה עם קישור חתימה מאובטח",
    "התראות על תגובות ותשובות חדשות עם תצוגה מקדימה",
    "התראת דחייה לבעל המסמך",
    "תור משימות אסינכרוני - שליחה אמינה גם בעומס",
]
add_text(s, Inches(7.1), Inches(2.05), Inches(5.5), Inches(0.5), "תגובות ואנוטציות", size=18, bold=True, color=BLUE)
add_bullets(s, Inches(6.9), Inches(2.6), Inches(5.7), Inches(4.0), left, size=13.5, space_after=10)
add_text(s, Inches(0.7), Inches(2.05), Inches(5.5), Inches(0.5), "התראות מייל", size=18, bold=True, color=BLUE)
add_bullets(s, Inches(0.5), Inches(2.6), Inches(5.7), Inches(4.0), right, size=13.5, space_after=10)
n += 1; add_footer(s, n, TOTAL)

# ---------- Slide 12: Audit & Security ----------
s = prs.slides.add_slide(blank)
add_bg(s)
section_header(s, "אבטחה ותאימות", "מעקב מלא ובקרת גישה")
cards = [
    ("יומן ביקורת מלא (Audit Trail)", "תיעוד כל פעולה: יצירה, העלאה, שינויי סטטוס, השלמת שלבים, חתימות, דחיות ודילוגים - עם זהות פעיל וחותמת זמן לכל אירוע"),
    ("היסטוריית חתימות", "שמירת כל חתימה עם זמן, זהות חותם, תמונת החתימה ומיקומה המדויק במסמך"),
    ("הרשאות וזיהוי משתמשים (Clerk)", "כניסה מאובטחת (SSO, Google, אימייל), תפקידי צוות (מנהל, חבר, אורח), ובידוד נתונים מלא בין לקוחות (Multi-tenant)"),
    ("בקרת גישה למסמכים", "רק בעל המסמך מנהל אותו; חותמים אורחים מקבלים טוקן מוגבל בזמן בלבד"),
]
top = Inches(2.0)
for i, (t, d) in enumerate(cards):
    feature_card(s, Inches(0.7), top + i*Inches(1.2), Inches(11.9), Inches(1.05), t, d, accent=NAVY)
n += 1; add_footer(s, n, TOTAL)

# ---------- Slide 13: Summary value table ----------
s = prs.slides.add_slide(blank)
add_bg(s)
section_header(s, "סיכום", "הערך העסקי בשורה אחת")
rows = [
    ("חיסכון בזמן", "מילוי אוטומטי ב-AI, חילוץ שדות, ושימוש חוזר בתבניות"),
    ("הפחתת שגיאות", "חילוץ נתונים אוטומטי, תבניות מוכנות מראש, ובדיקת שדות חובה"),
    ("תאימות משפטית ורגולטורית", "יומן ביקורת מקיף, אימות חתימות, ותמיכה בטופס הקנסות העירוני"),
    ("קלות שימוש", "אין צורך בהרשמה לחותמים אורחים, מיקום שדות חזותי, יצירת מסמך בלחיצה אחת"),
    ("עברית כשפת אם", "תמיכה מלאה ב-RTL, סיכום AI בעברית, וטפסים עירוניים ישראליים מובנים"),
]
top = Inches(2.0)
for i, (t, d) in enumerate(rows):
    feature_card(s, Inches(0.7), top + i*Inches(0.95), Inches(11.9), Inches(0.82), t, d, accent=GREEN)
n += 1; add_footer(s, n, TOTAL)

# ---------- Slide 14: Closing ----------
s = prs.slides.add_slide(blank)
add_bg(s, NAVY)
add_text(s, Inches(1), Inches(2.7), Inches(11.3), Inches(1.0),
          "DocFlows", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(3.8), Inches(10.3), Inches(0.8),
          "כל תהליך החוזה - מהעלאה ועד חתימה - במקום אחד, מהיר, מדויק ומתועד.",
          size=18, color=RGBColor(0xCF,0xDF,0xFF), align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(4.6), Inches(10.3), Inches(0.6),
          "תודה!", size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
n += 1; add_footer(s, n, TOTAL)

import os
out_path = os.path.join(os.path.dirname(__file__), "..", "DocFlows-Presentation-HE.pptx")
prs.save(out_path)
print("Saved:", os.path.abspath(out_path))
